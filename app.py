# ==========================================
# 1. INSTALL DEPENDENCIES
# ==========================================
# !pip install -q gradio groq python-docx python-pptx

import gradio as gr
from groq import Groq
import os
import json
from docx import Document
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ==========================================
# 2. SECURE API KEY RETRIEVAL
# ==========================================
import os
SECURE_API_KEY = os.getenv("GROQ_API_KEY")

# ==========================================
# 3. ENTERPRISE PRESENTATION GENERATOR
# ==========================================
def apply_heading_style(paragraph, text, is_main_title=False):
    paragraph.text = text
    run = paragraph.runs[0]
    run.font.bold = True
    run.font.name = 'Calibri'
    if is_main_title:
        run.font.size = Pt(22)
        run.font.color.rgb = RGBColor(0, 51, 102)
    else:
        run.font.size = Pt(16)
        run.font.color.rgb = RGBColor(0, 51, 102)

def apply_body_style(paragraph, text, is_bullet=False):
    if not is_bullet:
        paragraph.text = text
    else:
        paragraph.text = f"• {text}" 
        
    run = paragraph.runs[0]
    run.font.size = Pt(14)
    run.font.name = 'Calibri'
    run.font.color.rgb = RGBColor(64, 64, 64) 
    paragraph.space_after = Pt(12) 

def create_ppt(topic, slides_data, references):
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    background = title_slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 51, 102)
    title_shape = title_slide.shapes.title
    title_shape.text = f"Masterclass: {topic.upper()}"
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    title_shape.text_frame.paragraphs[0].font.bold = True
    subtitle = title_slide.placeholders[1]
    subtitle.text = "Comprehensive Professional Guide & Lecture Materials"
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(200, 200, 200)

    for slide_info in slides_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1]) 
        title_shape = slide.shapes.title
        title_shape.text = slide_info.get("title", "Topic Overview")
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
        
        body_shape = slide.placeholders[1]
        body_shape.left = Inches(0.5)
        body_shape.top = Inches(1.5)
        body_shape.width = Inches(9.0)
        body_shape.height = Inches(5.5)
        
        tf = body_shape.text_frame
        tf.clear() 
        tf.word_wrap = True

        apply_heading_style(tf.paragraphs[0], "Executive Summary", is_main_title=False)
        apply_body_style(tf.add_paragraph(), slide_info.get("executive_summary", ""))
        apply_heading_style(tf.add_paragraph(), "Detailed Analysis", is_main_title=False)
        apply_body_style(tf.add_paragraph(), slide_info.get("detailed_analysis", ""))
        apply_heading_style(tf.add_paragraph(), "Strategic Takeaways", is_main_title=False)
        for point in slide_info.get("key_takeaways", []):
            apply_body_style(tf.add_paragraph(), point, is_bullet=True)
        
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = "SPEAKER SCRIPT:\n\n" + slide_info.get("speaker_notes", "")

    filename = f"{topic.replace(' ', '_')}_Professional_Deck.pptx"
    prs.save(filename)
    return filename

def create_word_doc(topic, lesson, references, mcqs):
    doc = Document()
    doc.add_heading(f"Professional Guide: {topic}", 0)
    doc.add_heading("1. Comprehensive Masterclass", 1)
    doc.add_paragraph(lesson)
    doc.add_heading("2. Knowledge Assessment", 1)
    for i, q in enumerate(mcqs):
        doc.add_heading(f"Question {i+1}: {q['question']}", 2)
        for opt in q['options']:
            doc.add_paragraph(f"- {opt}")
        doc.add_paragraph(f"Answer: {q['answer']}").bold = True
        doc.add_paragraph("") 
    doc.add_heading("3. References & Citations", 1)
    doc.add_paragraph(references)
    
    filename = f"{topic.replace(' ', '_')}_Teachers_Guide.docx"
    doc.save(filename)
    return filename

# ==========================================
# 4. ROBUST AI PROMPT ENGINEERING
# ==========================================
def build_system_prompt(level: str) -> str:
    return f"""You are a top-tier management consultant. Audience Level: {level}. 
STRICT OUTPUT RULES:
1. Provide a professional masterclass with 5 sections.
2. IMPORTANT: Use ONLY single quotes (') for all text. NEVER use double quotes (") inside your JSON values.
3. Be concise to ensure the JSON structure is completed before the token limit.
You MUST return a VALID JSON object exactly like this:
{{
  "lesson_content": "Markdown text here...",
  "references_and_quotes": "List here...",
  "mcqs": [
    {{"question": "Q1", "options": ["A", "B", "C", "D"], "answer": "A"}}
  ],
  "ppt_slides": [
    {{
      "title": "Title", 
      "executive_summary": "Summary",
      "detailed_analysis": "Analysis",
      "key_takeaways": ["Point 1"],
      "speaker_notes": "Script"
    }}
  ]
}}"""

def generate_lesson(topic: str, level: str, progress=gr.Progress()):
    if not SECURE_API_KEY:
        return ("⚠️ API Key not found.", "", gr.update(), gr.update(), gr.update(), [], None, None)
    if not topic.strip():
        return ("⚠️ Please enter a topic.", "", gr.update(), gr.update(), gr.update(), [], None, None)

    try:
        # 0% to 20%: Initialization
        progress(0.05, desc="Initializing engine...")
        client = Groq(api_key=SECURE_API_KEY)
        
        progress(0.15, desc="Connecting to Groq API...")
        
        # 20% to 60%: The "Heavy Lifting" (Waiting for AI)
        progress(0.25, desc="Analyzing topic requirements...")
        progress(0.40, desc="Generating professional curriculum (this takes a moment)...")
        
        response = client.chat.completions.create(
            model="llama-3.1-8b-instant",
            messages=[
                {"role": "system", "content": build_system_prompt(level)},
                {"role": "user", "content": f"Create a definitive professional curriculum on: {topic}"}
            ],
            temperature=0.7,
            max_tokens=2500,
            response_format={"type": "json_object"} 
        )
        
        # 60% to 80%: Data Processing
        progress(0.65, desc="AI response received. Validating JSON structure...")
        raw_json = response.choices[0].message.content
        
        try:
            data = json.loads(raw_json)
        except json.JSONDecodeError:
            return ("⚠️ JSON Formatting Error. Please retry.", "", gr.update(visible=False), gr.update(visible=False), gr.update(visible=False), [], None, None)
            
        progress(0.75, desc="Extracting lesson modules and quiz data...")
        lesson = data.get("lesson_content", "")
        refs = data.get("references_and_quotes", "")
        mcqs = data.get("mcqs", [])
        slides = data.get("ppt_slides", [])

        # 80% to 100%: File Creation
        progress(0.85, desc="Generating PowerPoint slides (PPTX)...")
        ppt_file = create_ppt(topic, slides, refs)
        
        progress(0.95, desc="Finalizing Teacher's Guide (DOCX)...")
        word_file = create_word_doc(topic, lesson, refs, mcqs)

        q1, q2, q3 = mcqs[0], mcqs[1], mcqs[2]
        correct_answers = [q1["answer"], q2["answer"], q3["answer"]]

        progress(1.0, desc="Success! Files ready for download.")
        return (
            lesson, refs,
            gr.update(label=q1["question"], choices=q1["options"], visible=True),
            gr.update(label=q2["question"], choices=q2["options"], visible=True),
            gr.update(label=q3["question"], choices=q3["options"], visible=True),
            correct_answers, word_file, ppt_file         
        )

    except Exception as e:
        return (f"❌ API Error: {str(e)}", "", gr.update(), gr.update(), gr.update(), [], None, None)

def evaluate_quiz(ans1, ans2, ans3, correct_answers):
    if not correct_answers: return "Generate a lesson first!"
    user_answers = [ans1, ans2, ans3]
    score = 0
    feedback = "### Performance Review\n\n"
    for i, (user, correct) in enumerate(zip(user_answers, correct_answers)):
        if user == correct:
            score += 1
            feedback += f"✅ **Q{i+1}:** Correct!\n"
        else:
            feedback += f"❌ **Q{i+1}:** Incorrect. The correct answer is: *{correct}*\n"
    feedback += f"\n**Total Score: {score}/3**"
    return feedback

# ==========================================
# 5. UI GENERATION
# ==========================================
def create_app():
    with gr.Blocks(theme=gr.themes.Soft(primary_hue="slate")) as demo:
        gr.Markdown("# 🏢 Enterprise AI Masterclass Generator")
        gr.Markdown("Generates corporate curriculum. Optimized for API rate limits.")
        
        correct_answers_state = gr.State([])

        with gr.Row():
            with gr.Column(scale=1):
                topic_input = gr.Textbox(label="Target Topic", placeholder="e.g. Corporate Finance")
                level_input = gr.Dropdown(label="Audience Expertise", choices=["Beginner", "Intermediate", "Advanced"], value="Intermediate")
                submit_btn = gr.Button("Generate Professional Deck 🚀", variant="primary")
            
            with gr.Column(scale=2):
                with gr.Tabs():
                    with gr.Tab("📚 The Brief"):
                        lesson_display = gr.Markdown(value="*Curriculum will appear here...*")
                    with gr.Tab("📖 Citations"):
                        refs_display = gr.Markdown(value="*References will appear here...*")
                    with gr.Tab("📝 Knowledge Check"):
                        q1_radio = gr.Radio(label="Question 1", visible=False)
                        q2_radio = gr.Radio(label="Question 2", visible=False)
                        q3_radio = gr.Radio(label="Question 3", visible=False)
                        quiz_submit_btn = gr.Button("Evaluate Answers", variant="secondary")
                        quiz_results_display = gr.Markdown()
                    with gr.Tab("⬇️ Export Assets"):
                        word_download = gr.File(label="Teacher's Guide (.docx)")
                        ppt_download = gr.File(label="Professional Deck (.pptx)")

        submit_btn.click(
            fn=generate_lesson,
            inputs=[topic_input, level_input],
            outputs=[lesson_display, refs_display, q1_radio, q2_radio, q3_radio, correct_answers_state, word_download, ppt_download]
        )

        quiz_submit_btn.click(
            fn=evaluate_quiz,
            inputs=[q1_radio, q2_radio, q3_radio, correct_answers_state],
            outputs=[quiz_results_display]
        )
        
    return demo

if __name__ == "__main__":
    create_app().launch(debug=True)
